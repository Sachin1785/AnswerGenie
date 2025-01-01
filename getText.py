import pytesseract
from PIL import Image
import fitz  # PyMuPDF
import os
import docx
import pptx

# Path to tesseract executable (adjust it based on your installation)
pytesseract.pytesseract.tesseract_cmd = r'C:\Program Files\Tesseract-OCR\tesseract.exe'  # Update for your OS

def extract_text_from_image(image_path):
    """Extract text from a single image using pytesseract."""
    image = Image.open(image_path)
    return pytesseract.image_to_string(image)

def extract_text_from_pdf(pdf_path):
    """Extract text from a PDF using PyMuPDF (fitz) and fall back on OCR for images."""
    doc = fitz.open(pdf_path)
    text = ""

    for page_num in range(doc.page_count):
        page = doc.load_page(page_num)
        # Extract text from the page
        page_text = page.get_text("text")
        if page_text.strip():
            text += f"\n--- Page {page_num + 1} ---\n" + page_text.replace('\n', ' ')
        else:
            # If no text is found, fall back to OCR (for scanned pages)
            pix = page.get_pixmap()
            image = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
            image_text = pytesseract.image_to_string(image)
            text += f"\n--- Page {page_num + 1} (OCR) ---\n" + image_text.replace('\n', ' ')

    return text

def extract_text_from_txt(txt_path):
    """Extract text from a txt file."""
    with open(txt_path, 'r', encoding='utf-8') as file:
        return file.read()

def extract_text_from_docx(docx_path):
    """Extract text from a docx file."""
    doc = docx.Document(docx_path)
    return '\n'.join([para.text for para in doc.paragraphs])

def extract_text_from_pptx(pptx_path):
    """Extract text from a pptx file."""
    prs = pptx.Presentation(pptx_path)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return '\n'.join(text)

def save_text_to_file(text, output_path):
    """Save extracted text to a text file."""
    with open(output_path, "w", encoding="utf-8") as text_file:
        text_file.write(text)

def process_file(file_path):
    """Process input PDF, image, txt, docx, or pptx file."""
    if file_path.endswith('.pdf'):
        text = extract_text_from_pdf(file_path)
    elif file_path.lower().endswith(('.png', '.jpg', '.jpeg', '.bmp', '.tiff')):
        text = extract_text_from_image(file_path)
    elif file_path.lower().endswith('.txt'):
        text = extract_text_from_txt(file_path)
    elif file_path.lower().endswith('.docx'):
        text = extract_text_from_docx(file_path)
    elif file_path.lower().endswith('.pptx'):
        text = extract_text_from_pptx(file_path)
    else:
        raise ValueError("Unsupported file format. Please provide a PDF, image, txt, docx, or pptx file.")

    output_path = os.path.splitext(file_path)[0] + ".txt"
    save_text_to_file(text, output_path)
    print(f"Text extracted and saved to {output_path}")

# Example usage

# process_file(input("Enter the path to the PDF or image file: "))
process_file("UHV Chpt 1-15 all (1)[1].pdf")  # For testing purposes