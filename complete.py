import os
import faiss
import numpy as np
from sentence_transformers import SentenceTransformer
from groq import Groq
from dotenv import load_dotenv
import tkinter as tk
from tkinter import filedialog, messagebox, simpledialog
import pytesseract
from PIL import Image
import fitz  # PyMuPDF
import docx
import pptx

# Path to tesseract executable (adjust it based on your installation)
pytesseract.pytesseract.tesseract_cmd = r'C:\Program Files\Tesseract-OCR\tesseract.exe'  # Update for your OS

# Load environment variables
load_dotenv()

# Initialize the LLM client
api_key = os.getenv('GROQ_API_KEY')
client = Groq(api_key=api_key)
model_version = "llama3-8b-8192"
messages = [
    {"role": "system", "content": "You are a helpful assistant. Ready to answer any question with relevant answers"}
]

# Initialize the model and FAISS index
model = SentenceTransformer('all-MiniLM-L6-v2')
dimension = model.get_sentence_embedding_dimension()
index = faiss.IndexFlatL2(dimension)
sentences = []

# Function to extract text from files
def process_file(file_path):
    """Process input PDF, image, txt, docx, or pptx file."""
    def extract_text_from_image(image_path):
        image = Image.open(image_path)
        return pytesseract.image_to_string(image)

    def extract_text_from_pdf(pdf_path):
        doc = fitz.open(pdf_path)
        text = ""
        for page_num in range(doc.page_count):
            page = doc.load_page(page_num)
            page_text = page.get_text("text")
            if page_text.strip():
                text += f"\n--- Page {page_num + 1} ---\n" + page_text.replace('\n', ' ')
            else:
                # If no text is found, fall back to OCR (for scanned pages)
                pix = page.get_pixmap()
                image = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                image_text = pytesseract.image_to_string(image)
                text += f"\n--- Page {page_num + 1} (OCR) ---\n" + image_text.replace('\n', ' ')
        return text

    def extract_text_from_txt(txt_path):
        with open(txt_path, 'r', encoding='utf-8') as file:
            return file.read()

    def extract_text_from_docx(docx_path):
        doc = docx.Document(docx_path)
        return '\n'.join([para.text for para in doc.paragraphs])

    def extract_text_from_pptx(pptx_path):
        prs = pptx.Presentation(pptx_path)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return '\n'.join(text)

    def save_text_to_file(text, output_path):
        with open(output_path, "w", encoding="utf-8") as text_file:
            text_file.write(text)

    if file_path.lower().endswith('.pdf'):
        text = extract_text_from_pdf(file_path)
    elif file_path.lower().endswith(('.png', '.jpg', '.jpeg', '.bmp', '.tiff')):
        text = extract_text_from_image(file_path)
    elif file_path.lower().endswith('.txt'):
        text = extract_text_from_txt(file_path)
    elif file_path.lower().endswith('.docx'):
        text = extract_text_from_docx(file_path)
    elif file_path.lower().endswith('.pptx'):
        text = extract_text_from_pptx(file_path)
    else:
        raise ValueError("Unsupported file format.")
    output_path = os.path.splitext(file_path)[0] + ".txt"
    save_text_to_file(text, output_path)
    print(f"Text extracted and saved to {output_path}")

# Function to convert text to embeddings
def text_to_embeddings(text, model):
    sentences_list = text.split('\n')
    embeddings = model.encode(sentences_list)
    return embeddings, sentences_list

# Function to perform semantic search
def semantic_search(query, model, index, top_k=5):
    query_embedding = model.encode([query])
    D, I = index.search(query_embedding, top_k)
    return I[0]

# Function to add file to index
def add_file_to_index(file_path):
    # Extract text from the file
    process_file(file_path)
    # Load the extracted text
    text_file = os.path.splitext(file_path)[0] + ".txt"
    with open(text_file, 'r', encoding='utf-8') as file:
        text = file.read()
    # Convert text to embeddings
    embeddings, new_sentences = text_to_embeddings(text, model)
    # Add embeddings to the index
    index.add(embeddings)
    # Update the sentences list
    sentences.extend(new_sentences)

# Function to answer user questions
def answer_question(query):
    # Perform semantic search to retrieve relevant sentences
    top_k = 5
    results = semantic_search(query, model, index, top_k)
    retrieved_sentences = [sentences[idx] for idx in results if idx < len(sentences)]
    # Combine retrieved sentences
    context = "\n".join(retrieved_sentences)
    # Prepare prompt for the LLM
    prompt = f"Question: {query}\nContext: {context}\nAnswer:"
    # Get response from the LLM
    completion = client.chat.completions.create(
        model=model_version,
        messages=messages + [{"role": "user", "content": prompt}],
        temperature=0.85,
        top_p=1,
        stream=False,
        stop=None,
    )
    response = completion.choices[0].message.content.strip()
    # Update conversation history
    messages.append({"role": "user", "content": query})
    messages.append({"role": "assistant", "content": response})
    return response

def create_gui():
    # Initialize the main window with a background color
    root = tk.Tk()
    root.title("RAG Application")
    root.geometry("600x400")
    root.configure(bg='#2c3e50')  # Dark blue-gray background

    # Function to add files
    def add_files():
        file_paths = filedialog.askopenfilenames(title="Select files to upload")
        for file_path in file_paths:
            if os.path.exists(file_path):
                add_file_to_index(file_path)
        status_label.config(text="Files added to the vector database.", fg='white')

    # Create a frame for the file upload section with background color
    upload_frame = tk.Frame(root, bg='#2c3e50')
    upload_frame.pack(pady=10)

    add_button = tk.Button(upload_frame, text="Add Files", command=add_files, bg='#1abc9c', fg='white')
    add_button.pack()

    # Create a frame for the question and answer section with background color
    qa_frame = tk.Frame(root, bg='#2c3e50')
    qa_frame.pack(pady=10)

    question_label = tk.Label(qa_frame, text="Enter your question:", bg='#2c3e50', fg='white')
    question_label.pack(anchor='w')

    question_entry = tk.Entry(qa_frame, width=80, bg='#ecf0f1')
    question_entry.pack()

    # Function to handle question submission
    def ask_question():
        query = question_entry.get()
        if query:
            response = answer_question(query)
            answer_text.delete(1.0, tk.END)
            answer_text.insert(tk.END, response)
            question_entry.delete(0, tk.END)  # Clear the entry after submission

    ask_button = tk.Button(qa_frame, text="Ask", command=ask_question, bg='#3498db', fg='white')
    ask_button.pack(pady=5)

    answer_label = tk.Label(qa_frame, text="Answer:", bg='#2c3e50', fg='white')
    answer_label.pack(anchor='w')

    answer_text = tk.Text(qa_frame, height=10, width=80, bg='#ecf0f1')
    answer_text.pack()

    # Status label for feedback messages
    status_label = tk.Label(root, text="", bg='#2c3e50', fg='white')
    status_label.pack()

    # Start the GUI event loop
    root.mainloop()

# Replace the main loop with the GUI
if __name__ == "__main__":
    create_gui()
